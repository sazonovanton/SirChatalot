#!/usr/bin/env python3
'''
Class for converting PDF, DOC, DOCX, PPT, PPTX, TXT files to just a plain text.
Class for processing text files and inserting them into the RAG database.
'''
import asyncio
import os

# File processing
import aiofiles
import fitz  # PyMuPDF
import docx 
import pptx  

# RAG
import chromadb
from chromadb.config import Settings
from chromadb.utils import embedding_functions
from chatutils.emb_engines import get_embeddings_engine
import uuid
import re

# Configuring
import configparser
config = configparser.ConfigParser()
config.read('./data/.config', encoding='utf-8')
LogLevel = config.get("Logging", "LogLevel") if config.has_option("Logging", "LogLevel") else "WARNING"

# Logging
import logging
from logging.handlers import TimedRotatingFileHandler
logger = logging.getLogger("SirChatalot-FilesProc")
LogLevel = getattr(logging, LogLevel.upper())
logger.setLevel(LogLevel)
handler = TimedRotatingFileHandler('./logs/sirchatalot.log',
                                       when="D",
                                       interval=1,
                                       backupCount=7,
                                       encoding='utf-8')
handler.setFormatter(logging.Formatter('%(name)s - %(asctime)s - %(levelname)s - %(message)s',"%Y-%m-%d %H:%M:%S"))
logger.addHandler(handler)


######### Files Processor #########
class FilesProcessor:
    def __init__(self, files_path: str = "./data/files") -> None:
        self.path = files_path
        if not os.path.exists(self.path):
            os.makedirs(self.path, exist_ok=True)
    
    async def convert_to_text(self, filename: str) -> str:
        """
        Converts a file to plain text.

        Parameters:
            filename: name of the file to convert
        
        Returns the text content of the file.
        """
        try:
            text = ""
            if filename.endswith(".pdf"):
                text = await self.convert_pdf_to_text(filename)
            elif filename.endswith(".docx"):
                text = await self.convert_docx_to_text(filename)
            elif filename.endswith(".doc"):
                text = await self.convert_doc_to_text(filename)
            elif filename.endswith(".pptx"):
                text = await self.convert_pptx_to_text(filename)
            elif filename.endswith(".ppt"):
                text = await self.convert_ppt_to_text(filename)
            elif filename.endswith(".txt") or filename.endswith(".md") or filename.endswith(".csv") or filename.endswith(".log"):
                text = await self.read_text_file(filename)
            return text
        except KeyboardInterrupt:
            logger.error("Conversion cancelled by user.")
            raise KeyboardInterrupt
        except Exception as e:
            logger.error(f"Error converting file to text: {e}")
            return None
    
    async def convert_pdf_to_text(self, filename: str) -> str:
        """
        Converts a PDF file to text using PyMuPDF.
        """
        # # full_path = os.path.join(self.path, filename)
        full_path = filename
        try:
            text = await self._extract_pdf_text(full_path)
            return text
        except Exception as e:
            logger.error(f"Error converting PDF to text: {e}")
            return None
    
    async def _extract_pdf_text(self, file_path: str) -> str:
        """
        Extracts text from a PDF file using PyMuPDF in an async-friendly way.
        Processing pages one by one to avoid blocking the event loop for too long.
        """
        # Open the document in a separate thread
        doc = await asyncio.to_thread(fitz.open, file_path)
        
        text_parts = []
        try:
            # Process each page asynchronously
            for page_num in range(len(doc)):
                # Extract text from page - this could be CPU intensive, so run in thread
                page = doc[page_num]  # Get page object
                page_text = await asyncio.to_thread(page.get_text)
                text_parts.append(page_text)
                
                # Yield control periodically to not block the event loop
                if page_num % 5 == 0:  # Every 5 pages
                    await asyncio.sleep(0)

        except Exception as e:
            logger.error(f"Error processing page {page_num}: {e}")
        finally:
            # Make sure we close the document
            await asyncio.to_thread(doc.close)
            return "\n".join(text_parts)

        
    async def convert_docx_to_text(self, filename: str) -> str:
        """
        Converts DOCX file to text.
        """
        # full_path = os.path.join(self.path, filename)
        full_path = filename
        try:
            text = await asyncio.to_thread(self._extract_docx, full_path)
            return text
        except Exception as e:
            logger.error(f"Error converting DOCX to text: {e}")
            return None
        
    def _extract_docx(self, file_path: str) -> str:
        """
        Extracts text from a DOCX file using python-docx.
        """
        doc = docx.Document(file_path)
        paragraphs = [para.text for para in doc.paragraphs]
        return "\n".join(paragraphs)

    async def convert_doc_to_text(self, filename: str) -> str:
        """
        Converts a DOC file to text using catdoc.
        """
        # full_path = os.path.join(self.path, filename)
        full_path = filename
        try:
            # Run catdoc asynchronously
            proc = await asyncio.create_subprocess_exec(
                "catdoc", full_path,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE
            )
            stdout, stderr = await proc.communicate()
            
            if stderr:
                logger.warning(f"Warning while converting DOC file: {stderr.decode()}")
            
            return stdout.decode()
        except Exception as e:
            logger.error(f"Error converting DOC to text: {e}")
            return None

    async def convert_pptx_to_text(self, filename: str) -> str:
        """
        Converts PPTX file to text
        """
        # full_path = os.path.join(self.path, filename)
        full_path = filename
        try:
            text = await asyncio.to_thread(self._extract_pptx_text, full_path)
            return text
        except Exception as e:
            logger.error(f"Error converting PPTX to text: {e}")
            return None

    def _extract_pptx_text(self, full_path: str) -> str:
        """
        Extracts text from a PPTX file using python-pptx.
        """
        prs = pptx.Presentation(full_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
        return "\n".join(texts)

    async def convert_ppt_to_text(self, filename: str) -> str:
        """
        Converts a PPT file to text using catppt.
        """
        # full_path = os.path.join(self.path, filename)
        full_path = filename
        try:
            # Run catppt asynchronously
            proc = await asyncio.create_subprocess_exec(
                "catppt", full_path,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE
            )
            stdout, stderr = await proc.communicate()
            
            if stderr:
                logger.warning(f"Warning while converting PPT file: {stderr.decode()}")
            
            return stdout.decode()
        except Exception as e:
            logger.error(f"Error converting PPT to text: {e}")
            return None

    async def read_text_file(self, filename: str) -> str:
        """
        Reads text from a text file.
        """
        # full_path = os.path.join(self.path, filename)
        full_path = filename
        try:
            async with aiofiles.open(full_path, mode="r", encoding="utf-8") as f:
                content = await f.read()
            return content
        except Exception as e:
            logger.error(f"Error reading text file: {e}")
            return None
        

######### RAG #########
class FilesRAG:
    def __init__(self, chromadb_path: str = "./data/files/chromadb", files_path: str = "./data/files") -> None:
        self.path = files_path
        if not os.path.exists(self.path):
            os.makedirs(self.path, exist_ok=True)

        self.emb_engine = get_embeddings_engine()
        
        self.client = chromadb.PersistentClient(
            path=chromadb_path,
            settings=Settings()
        )
        
        self.init_collection()

    def init_collection(self) -> None:
        """
        Initializes the collection for storing texts and vectors.
        """
        try:
            self.collection = self.client.get_collection(
                name="files",
                embedding_function=None
            )
            logger.info("Collection 'files' found. Using existing collection.")
        except KeyboardInterrupt:
            logger.error("Initialization cancelled by user.")
            raise KeyboardInterrupt
        except Exception as e:
            logger.info(f"Collection 'files' not found. Creating new collection. Error: {e}")
            self.collection = self.client.create_collection(
                name="files",
                embedding_function=None
            )

    async def insert_texts(self, texts, metadata) -> bool:
        """
        Inserts texts into the collection. Calculates embeddings for the texts using the embeddings engine.

        Parameters:
            texts: list of texts (for example, list of strings)
            metadata: list of metadata (for example, list of dictionaries)
                we should have user_id in metadata: 
                [
                    {"user_id": "user1", "filename": "doc1.txt"},
                ]

        Returns True if texts were successfully inserted, False otherwise.
        """
        try:
            vectors, _ = await self.emb_engine.get_embeddings(texts)
            ids = [str(uuid.uuid4()) for _ in texts]  # Generate a unique ID for each text
            self.collection.add(
                ids=ids,
                embeddings=vectors,
                documents=texts,
                metadatas=metadata
            )
            logger.info("Texts inserted successfully.")
            return True
        except KeyboardInterrupt:
            logger.error("Insertion cancelled by user.")
            raise KeyboardInterrupt
        except Exception as e:
            logger.error(f"Error inserting texts: {e}")
            return False

    async def search_text(self, text, user_id, n_results=4, filter: dict = None, max_distance=1.0) -> list:
        """
        Searches for similar texts based on the given vector. Embeddings are calculated for the text using the embeddings engine.

        Parameters:
            text: text to search for
            user_id: user id to filter results (also searches for user_id = 'common')
            n_results: number of results to return
            filter: dictionary with filter conditions
            max_distance: maximum distance for search
        
        Returns a list of results.
        """
        try:
            vector, _ = await self.emb_engine.get_embeddings(text)

            where_clause = filter.copy() if filter else {}

            # Query for user-specific content
            where_clause["user_id"] = user_id
            result_user = self.collection.query(
                query_embeddings=[vector],
                n_results=n_results,
                where=where_clause
            )

            # Query for common content
            where_clause["user_id"] = "common"
            result_common = self.collection.query(
                query_embeddings=[vector],
                n_results=n_results,
                where=where_clause
            )

            # Initialize with the user results
            result = {}
            
            # Function to safely initialize or extend lists
            def init_or_extend(target_dict, key, source_dict):
                if key in source_dict and source_dict[key] and len(source_dict[key]) > 0:
                    if key not in target_dict or not target_dict[key]:
                        target_dict[key] = source_dict[key]
                    else:
                        target_dict[key].extend(source_dict[key])

            # Process user results
            if result_user and all(key in result_user for key in ["ids", "embeddings", "documents", "metadatas", "distances"]):
                for key in ["ids", "embeddings", "documents", "metadatas", "distances"]:
                    init_or_extend(result, key, result_user)

            # Process common results
            if result_common and all(key in result_common for key in ["ids", "embeddings", "documents", "metadatas", "distances"]):
                for key in ["ids", "embeddings", "documents", "metadatas", "distances"]:
                    init_or_extend(result, key, result_common)

            if not result or "ids" not in result or not result["ids"]:
                logger.info("No results found for the given search.")
                return []
            
            # Sort results based on distance
            if all(key in result for key in ["ids", "embeddings", "documents", "metadatas", "distances"]):
                # Create a list of tuples for sorting
                combined = list(zip(result["ids"], result["embeddings"], result["documents"], 
                                  result["metadatas"], result["distances"]))
                # Sort by distance (5th element, index 4)
                combined.sort(key=lambda x: x[4])
                
                # Unpack the sorted results
                result["ids"], result["embeddings"], result["documents"], result["metadatas"], result["distances"] = \
                    [list(item) for item in zip(*combined)]

            logger.debug(f"Search results: {result}")
            return result
        except KeyboardInterrupt:
            logger.error("Search cancelled by user.")
            raise KeyboardInterrupt
        except Exception as e:
            logger.error(f"Error searching texts: {e}")
            return None
        
    async def semantic_search(self, text, user_id, n_results=4, filter: dict = None, max_distance=1.0) -> list:
        """
        Searches for similar texts based on the given text. Embeddings are calculated for the text using the embeddings engine.

        Parameters:
            text: text to search for
            user_id: user id to filter results
            n_results: number of results to return
            filter: dictionary with filter conditions
            max_distance: maximum distance for search

        Returns a list of results - formatted to {"filename": [text, ...]}
        Uses search_text to get the results and then formats them based on the filename.
        """
        try:
            if n_results is None:
                n_results = 4
            logger.debug(f"Searching for similar texts based on: {text} (for user: {user_id})")
            results = await self.search_text(text, user_id, n_results, filter, max_distance)
            if results:
                formatted_results = {}
                for i, result in enumerate(results["metadatas"]):
                    if result and "filename" in result[0]:
                        filename = result[0]["filename"]
                        if filename not in formatted_results:
                            formatted_results[filename] = []
                        formatted_results[filename].extend(results["documents"][i])  
                logger.debug(f"Semantic search results: {formatted_results}")
                return formatted_results
            else:
                return {}
        except KeyboardInterrupt:
            logger.error("Semantic search cancelled by user.")
            raise KeyboardInterrupt
        except Exception as e:
            logger.error(f"Error performing semantic search: {e}")
            return []

    async def remove_texts(self, filename: str) -> bool:
        """
        Removes texts from the collection based on the filename.

        Parameters:
            filename: name of the file to remove
        
        Returns True if texts were successfully removed, False otherwise.
        """
        try:
            self.collection.delete(where={"filename": filename})  
            logger.info("Texts removed successfully.")
            return True
        except Exception as e:
            logger.error(f"Error removing texts: {e}")
            return False

    async def remove_text_user(self, user_id) -> bool:
        """
        Removes texts from the collection based on the user_id.

        Parameters:
            user_id: user id to remove
        
        Returns True if texts were successfully removed, False otherwise.
        """
        try:
            self.collection.delete(where={"user_id": user_id})
            logger.info("Texts removed successfully.")
            return True
        except Exception as e:
            logger.error(f"Error removing texts: {e}")
            return False
        
    async def user_files(self, user_id) -> list:
        """
        Get all the files associated with a user.

        Parameters:
            user_id: User ID to filter files
        
        Returns a list of unique filenames associated with the user.
        """
        try:
            # Use the get method instead of query to retrieve entries by metadata
            results = self.collection.get(
                where={"user_id": user_id},
                include=["metadatas"]
            )
            
            # Extract unique filenames from the metadata
            filenames = set()
            if results and "metadatas" in results and results["metadatas"]:
                for metadata in results["metadatas"]:
                    if metadata and "filename" in metadata:
                        filenames.add(metadata["filename"])
            
            return list(filenames)
        except KeyboardInterrupt:
            logger.error("User files retrieval cancelled by user.")
            raise KeyboardInterrupt
        except Exception as e:
            logger.error(f"Error retrieving user files: {e}")
            return []

    async def process_text(self, text: str, user_id, filename: str, chunk_size: int = 600, overlap_percent: float = 0.1) -> bool:
        """
        Process a text by intelligently splitting it into overlapping chunks at natural boundaries
        and inserting them into the collection.
        
        Parameters:
            text: The text content to process
            user_id: User ID to associate with the text chunks
            filename: Name of the file to associate with the text chunks
            chunk_size: Target size of each text chunk in characters
            overlap_percent: Percentage of overlap between chunks (0.0 to 1.0)
            
        Returns:
            True if the text was successfully processed and inserted, False otherwise
        """
        try:
            # Handle empty text case
            if not text or len(text.strip()) == 0:
                logger.warning(f"Empty text provided for file: {filename}")
                return False
                
            # Calculate the overlap size in characters
            overlap_size = int(chunk_size * overlap_percent)
            
            # Ensure we have a minimum chunk size after accounting for overlap
            effective_chunk_size = chunk_size - overlap_size
            if effective_chunk_size <= 0:
                logger.error(f"Invalid chunk configuration: chunk_size = {chunk_size}, overlap_size = {overlap_size}")
                return False
            
            # Split the text into paragraphs first
            paragraphs = [p.strip() for p in re.split(r'\n\s*\n', text) if p.strip()]
            
            # Initialize for collecting chunks
            chunks = []
            metadata = []
            current_chunk = ""
            current_start = 0
            chunk_number = 1
            
            for paragraph in paragraphs:
                # If adding this paragraph would exceed chunk size and we already have content, finalize the current chunk
                if len(current_chunk) > 0 and (len(current_chunk) + len(paragraph) + 2) > chunk_size:
                    chunks.append(current_chunk)
                    chunk_end = current_start + len(current_chunk)
                    metadata.append({
                        "user_id": user_id,
                        "filename": filename,
                        "chunk_number": chunk_number,
                        "start_char": current_start,
                        "end_char": chunk_end
                    })
                    
                    # Start a new chunk with overlap
                    if overlap_size > 0:
                        # Find a good breaking point for the overlap
                        overlap_text = current_chunk[-min(overlap_size, len(current_chunk)):]
                        
                        # Try to find a sentence boundary in the overlap region
                        sentence_match = re.search(r'[.!?]\s+[A-Z]', overlap_text)
                        if sentence_match:
                            overlap_pos = min(overlap_size, len(current_chunk)) - sentence_match.start()
                            current_chunk = current_chunk[-min(overlap_pos, len(current_chunk)):]
                        # If no sentence boundary, try to find a word boundary
                        else:
                            word_match = re.search(r'\s+\S+\s*$', overlap_text)
                            if word_match:
                                overlap_pos = min(overlap_size, len(current_chunk)) - word_match.start()
                                current_chunk = current_chunk[-min(overlap_pos, len(current_chunk)):]
                            else:
                                # Fall back to simple overlap
                                current_chunk = current_chunk[-min(overlap_size, len(current_chunk)):]
                        
                        current_start = chunk_end - len(current_chunk)
                    else:
                        current_chunk = ""
                        current_start = chunk_end
                    
                    chunk_number += 1
                
                # Handle paragraphs that exceed chunk_size on their own
                if len(paragraph) > chunk_size:
                    # If we have anything in the current chunk, finalize it first
                    if current_chunk:
                        chunks.append(current_chunk)
                        chunk_end = current_start + len(current_chunk)
                        metadata.append({
                            "user_id": user_id,
                            "filename": filename,
                            "chunk_number": chunk_number,
                            "start_char": current_start,
                            "end_char": chunk_end
                        })
                        chunk_number += 1
                        current_chunk = ""
                        current_start = chunk_end
                    
                    # Split long paragraph into multiple chunks
                    words = paragraph.split()
                    para_chunk = ""
                    for word in words:
                        if len(para_chunk) + len(word) + 1 > chunk_size:
                            chunks.append(para_chunk)
                            chunk_end = current_start + len(para_chunk)
                            metadata.append({
                                "user_id": user_id,
                                "filename": filename,
                                "chunk_number": chunk_number,
                                "start_char": current_start,
                                "end_char": chunk_end
                            })
                            
                            if overlap_size > 0:
                                overlap_text = para_chunk[-min(overlap_size, len(para_chunk)):]
                                sentence_match = re.search(r'[.!?]\s+[A-Z]', overlap_text)
                                if sentence_match:
                                    overlap_pos = min(overlap_size, len(para_chunk)) - sentence_match.start()
                                    para_chunk = para_chunk[-min(overlap_pos, len(para_chunk)):]
                                else:
                                    word_match = re.search(r'\s+\S+\s*$', overlap_text)
                                    if word_match:
                                        overlap_pos = min(overlap_size, len(para_chunk)) - word_match.start()
                                        para_chunk = para_chunk[-min(overlap_pos, len(para_chunk)):]
                                    else:
                                        para_chunk = para_chunk[-min(overlap_size, len(para_chunk)):]
                                
                                current_start = chunk_end - len(para_chunk)
                            else:
                                para_chunk = ""
                                current_start = chunk_end
                            
                            chunk_number += 1
                        
                        if para_chunk:
                            para_chunk += " " + word
                        else:
                            para_chunk = word
                    
                    if para_chunk:
                        current_chunk = para_chunk
                    else:
                        current_chunk = ""
                else:
                    # Add paragraph to current chunk
                    if current_chunk:
                        current_chunk += "\n\n" + paragraph
                    else:
                        current_chunk = paragraph
            
            # Don't forget the last chunk
            if current_chunk:
                chunks.append(current_chunk)
                chunk_end = current_start + len(current_chunk)
                metadata.append({
                    "user_id": user_id,
                    "filename": filename,
                    "chunk_number": chunk_number,
                    "start_char": current_start,
                    "end_char": chunk_end
                })
            
            # Insert the chunks into the collection
            if chunks:
                result = await self.insert_texts(chunks, metadata)
                logger.info(f"Processed text from {filename} into {len(chunks)} chunks using natural boundaries")
                return result
            else:
                logger.warning(f"No valid chunks generated from text in file: {filename}")
                return False
        
        except KeyboardInterrupt:
            logger.error("Text processing cancelled by user.")
            raise KeyboardInterrupt
        except Exception as e:
            logger.error(f"Error processing text: {e}")
            return False
            
        
if __name__ == "__main__":
    async def run_tests():
        # # Test for FilesProcessor
        # print("Starting tests for FilesProcessor...")

        # processor = FilesProcessor()

        # # Test conversion of a text file
        # text_content = await processor.read_text_file('test/sample.txt')
        # print(f"- Text file content: {text_content[:100]}...")  # Print first 100 characters

        # # Test conversion of a PDF file
        # pdf_text = await processor.convert_pdf_to_text('test/sample.pdf')
        # print(f"- PDF file content: {pdf_text[:100]}...")  # Print first 100 characters

        # # Test conversion of a DOCX file
        # docx_text = await processor.convert_docx_to_text('test/sample.docx')
        # print(f"- DOCX file content: {docx_text[:100]}...")  # Print first 100 characters

        # # Test conversion of a PPTX file
        # pptx_text = await processor.convert_pptx_to_text('test/sample.pptx')
        # print(f"- PPTX file content: {pptx_text[:100]}...")  # Print first 100 characters

        # # # Test conversion of a DOC file
        # # doc_text = await processor.convert_doc_to_text('test/sample.doc')
        # # print(f"- DOC file content: {doc_text[:100]}...")

        # # # Test conversion of a PPT file
        # # ppt_text = await processor.convert_ppt_to_text('test/sample.ppt')
        # # print(f"- PPT file content: {ppt_text[:100]}...")

        # print("Tests for FilesProcessor completed.")

        # Test for FilesRAG
        print("Starting tests for FilesRAG...")
        rag = FilesRAG()

        # Dummy data for testing insertion
        texts = [
            "The pineapple is a tropical plant with an edible fruit", 
            "A car, or an automobile, is a motor vehicle with wheels.", 
            "Paris is the capital and largest city of France.",
            "The peach is a deciduous tree first domesticated and cultivated in China.",
            "A bicycle, also called a bike or cycle, is a human-powered or motor-powered, pedal-driven, single-track vehicle, having two wheels attached to a frame, one behind the other.",
            "An airplane, an aeroplane, or a plane, is a fixed-wing aircraft that is propelled forward by thrust from a jet engine, propeller, or rocket engine."
        ]
        metadata = [
            {"user_id": "user1", "filename": "pineapple.txt"},
            {"user_id": "user1", "filename": "car.txt"},
            {"user_id": "user1", "filename": "paris.txt"},
            {"user_id": "user1", "filename": "peach.txt"},
            {"user_id": "user2", "filename": "bicycle.txt"},
            {"user_id": "user1", "filename": "airplane.txt"}
        ]
        
        # Test insertion of texts
        success = await rag.insert_texts(texts, metadata)
        print(f"- Insertion success: {success}")

        # Test searching texts
        search_result = await rag.search_text("car", "user1", n_results=2)
        print(f"- Search result: {search_result}")

        # Test getting user files
        user_files = await rag.user_files("user1")
        print(f"- User files: {user_files}")

        # Test semantic search
        semantic_result = await rag.semantic_search("fruit", "user1", n_results=2)
        print(f"- Semantic search result: {semantic_result}")

        # Test removal of texts by filename
        remove_success = await rag.remove_texts("pineapple.txt")
        print(f"- Remove by filename success: {remove_success}")

        # Test removal of texts by user_id
        remove_user_success = await rag.remove_text_user("user1")
        print(f"- Remove by user_id success: {remove_user_success}")

        print("Tests for FilesRAG completed.")

        # # Whole text processing
        # print("Starting whole text process")
        # rag = FilesRAG()
        # processor = FilesProcessor()

        # # Processing of a pdf file
        # pdf_text = await processor.convert_pdf_to_text('test/user_manual.pdf')
        # print(f"- PDF file content: {pdf_text[:100]}...")

        # # Process the text and insert into the collection
        # success = await rag.process_text(pdf_text, "user1", "user_manual.pdf")
        # print(f"- Text processing success: {success}")

        # # Find similar texts
        # search_result = await rag.search_text("power supply requirements", "user1", n_results=2)
        # print(f"- Search result: {search_result}")

        # print("Whole text process completed.")

    asyncio.run(run_tests())